"""
Document ingestion pipeline:
  1. Extract text (PDF → PyMuPDF, PPTX → python-pptx, DOCX → python-docx, else plaintext)
  2. Chunk text into overlapping passages
  3. Embed chunks via OpenAI text-embedding-3-small
  4. Store embeddings in ChromaDB
  5. Store chunk metadata in MySQL (ai_tutor_vector_chunks)
"""

import logging
from typing import List

from sqlalchemy.orm import Session

from app.core.config import settings

logger = logging.getLogger(__name__)

CHUNK_WORDS   = 400   # approximate words per chunk
CHUNK_OVERLAP = 50    # words of overlap between consecutive chunks


# Text extraction

def _extract_text(storage_path: str, mime_type: str) -> str:
    mime = (mime_type or "").lower()

    try:
        if "pdf" in mime or storage_path.endswith(".pdf"):
            import fitz  # PyMuPDF
            doc = fitz.open(storage_path)
            return "\n".join(page.get_text() for page in doc)

        if "presentationml" in mime or storage_path.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(storage_path)
            texts: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)

        if "wordprocessingml" in mime or storage_path.endswith(".docx"):
            import docx
            document = docx.Document(storage_path)
            return "\n".join(p.text for p in document.paragraphs if p.text.strip())

        # Plain text / markdown / CSV / fallback
        with open(storage_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    except Exception as e:
        logger.error(f"Text extraction failed for {storage_path}: {e}")
        return ""


# Chunking

def _chunk_text(text: str) -> List[str]:
    words = text.split()
    chunks: List[str] = []
    start = 0
    while start < len(words):
        end   = min(start + CHUNK_WORDS, len(words))
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        start += CHUNK_WORDS - CHUNK_OVERLAP
    return chunks


# Embeddings

def _get_embeddings(texts: List[str]) -> List[List[float]]:
    api_key = settings.OPENAI_API_KEY
    if not api_key:
        raise RuntimeError(
            "OPENAI_API_KEY is not set in .env. "
            "The AI Tutor feature requires an OpenAI API key."
        )
    import openai
    client = openai.OpenAI(api_key=api_key)

    all_embeddings: List[List[float]] = []
    batch_size = 100
    for i in range(0, len(texts), batch_size):
        batch = texts[i : i + batch_size]
        resp  = client.embeddings.create(model="text-embedding-3-small", input=batch)
        all_embeddings.extend(item.embedding for item in resp.data)
    return all_embeddings


# Public API

def ingest_document(doc, db: Session) -> int:
    """
    Parse, chunk, embed, and store a document. Returns the number of chunks created.
    Marks doc.is_indexed = 1 on success.
    """
    from app.models.ai_tutor import AiTutorVectorChunk
    from app.services.ai_tutor import vector_index_manager as vim

    text = _extract_text(doc.storage_path, doc.mime_type or "")
    if not text.strip():
        logger.warning(f"Document {doc.id} ({doc.original_filename}) yielded no text.")
        return 0

    chunks = _chunk_text(text)
    if not chunks:
        return 0

    embeddings = _get_embeddings(chunks)

    ids = [f"doc_{doc.id}_chunk_{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "tutor_id":    str(doc.tutor_id),
            "document_id": str(doc.id),
            "chapter_id":  str(doc.chapter_id or ""),
            "doc_type":    str(doc.doc_type.value if doc.doc_type else "other"),
            "filename":    doc.original_filename or "",
            "chunk_index": str(i),
        }
        for i in range(len(chunks))
    ]

    vim.add_chunks(doc.tutor_id, chunks, embeddings, metadatas, ids)

    # Refresh MySQL metadata (delete old chunks first to allow re-indexing)
    db.query(AiTutorVectorChunk).filter(AiTutorVectorChunk.document_id == doc.id).delete()
    for i, (chunk, vid) in enumerate(zip(chunks, ids)):
        db.add(AiTutorVectorChunk(
            tutor_id=doc.tutor_id,
            document_id=doc.id,
            chunk_index=i,
            chunk_text=chunk[:500],
            vector_id=vid,
        ))

    doc.is_indexed = 1
    db.commit()
    logger.info(f"Indexed document {doc.id} ({doc.original_filename}): {len(chunks)} chunks.")
    return len(chunks)


def remove_document_chunks(doc, db: Session) -> None:
    """Remove a document's vectors from ChromaDB and its metadata from MySQL."""
    from app.models.ai_tutor import AiTutorVectorChunk
    from app.services.ai_tutor import vector_index_manager as vim

    chunks = db.query(AiTutorVectorChunk).filter(AiTutorVectorChunk.document_id == doc.id).all()
    chunk_ids = [c.vector_id for c in chunks if c.vector_id]
    if chunk_ids:
        try:
            vim.delete_chunks_by_ids(doc.tutor_id, chunk_ids)
        except Exception as e:
            logger.warning(f"Failed to delete vector chunks for doc {doc.id}: {e}")
    db.query(AiTutorVectorChunk).filter(AiTutorVectorChunk.document_id == doc.id).delete()
    db.commit()


def ingest_transcript(transcript, db: Session) -> int:
    """Ingest an approved transcript's text into the vector store."""
    from app.models.ai_tutor import AiTutorVectorChunk
    from app.services.ai_tutor import vector_index_manager as vim

    text = transcript.approved_transcript or transcript.raw_transcript or ""
    if not text.strip():
        return 0

    chunks = _chunk_text(text)
    if not chunks:
        return 0

    embeddings = _get_embeddings(chunks)
    ids = [f"transcript_{transcript.id}_chunk_{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "tutor_id":      str(transcript.tutor_id),
            "transcript_id": str(transcript.id),
            "chapter_id":    str(transcript.chapter_id or ""),
            "doc_type":      "transcript",
            "filename":      f"Transcript #{transcript.id}",
            "chunk_index":   str(i),
        }
        for i in range(len(chunks))
    ]

    vim.add_chunks(transcript.tutor_id, chunks, embeddings, metadatas, ids)

    db.query(AiTutorVectorChunk).filter(AiTutorVectorChunk.transcript_id == transcript.id).delete()
    for i, (chunk, vid) in enumerate(zip(chunks, ids)):
        db.add(AiTutorVectorChunk(
            tutor_id=transcript.tutor_id,
            transcript_id=transcript.id,
            chunk_index=i,
            chunk_text=chunk[:500],
            vector_id=vid,
        ))

    transcript.is_indexed = 1
    db.commit()
    logger.info(f"Indexed transcript {transcript.id}: {len(chunks)} chunks.")
    return len(chunks)
